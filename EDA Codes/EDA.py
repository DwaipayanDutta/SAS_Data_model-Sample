import pandas as pd
import requests
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from scipy import stats

class DataAnalyzer:
    def __init__(self, url):
        self.url = url
        self.df = self.load_data()
        self.target_variable = None

    def load_data(self):
        # Ensure necessary libraries are imported in the executed code
        exec_code = """
import numpy as np  # Ensure NumPy is imported
# Your data loading logic here, assuming df is defined in this context.
"""
        exec(requests.get(self.url).text + exec_code, globals())
        return globals()['df']  # Assuming df is defined in the executed script

    def set_target_variable(self, target_variable):
        self.target_variable = target_variable

    def get_continuous_vars(self):
        return self.df.select_dtypes(include=['float64', 'int64']).columns

    def get_categorical_vars(self):
        return self.df.select_dtypes(include=['object']).columns


class Plotter:
    def __init__(self, presentation):
        self.ppt = presentation

    def set_title_format(self, title_shape):
        title_shape.text_frame.paragraphs[0].font.name = 'Calibri'
        title_shape.text_frame.paragraphs[0].font.size = Pt(16)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color

    def save_plot_to_ppt(self, plot_func, slide_title):
        plt.figure(figsize=(10, 5))
        plot_func()
        plt.title(slide_title)
        plt.tight_layout()

        # Save the plot to a temporary file
        plt.savefig('temp_plot.png')
        plt.close()

        # Add a slide to the presentation
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[5])  # Title Only layout
        title = slide.shapes.title
        title.text = slide_title.title()  # Convert to title case

        self.set_title_format(title)  # Set title formatting

        # Add image to slide
        left = Inches(1)
        top = Inches(1.5)
        slide.shapes.add_picture('temp_plot.png', left, top, width=Inches(8))


class EDA:
    def __init__(self, analyzer, plotter):
        self.analyzer = analyzer
        self.plotter = plotter

    def perform_kde_analysis(self):
        continuous_vars = self.analyzer.get_continuous_vars()
        
        for var in continuous_vars:
            self.plotter.save_plot_to_ppt(
                lambda: sns.kdeplot(data=self.analyzer.df, x=var, hue=self.analyzer.target_variable, fill=True, common_norm=False),
                f'KDE Plot of {var} by {self.analyzer.target_variable}'
            )

    def perform_outlier_analysis(self):
        continuous_vars = self.analyzer.get_continuous_vars()
        
        z_scores = stats.zscore(self.analyzer.df[continuous_vars])
        abs_z_scores = abs(z_scores)
        outliers = (abs_z_scores > 3).any(axis=1)

        # Add outlier information to DataFrame
        self.analyzer.df['Outlier'] = outliers
        
        # Visualize Outliers with Box Plots for Continuous Variables
        for var in continuous_vars:
            self.plotter.save_plot_to_ppt(
                lambda: sns.boxplot(x='Outlier', y=var, data=self.analyzer.df), 
                f'Box Plot of {var} Showing Outliers'
            )

    def perform_bivariate_analysis(self):
        continuous_vars = self.analyzer.get_continuous_vars()
        
        for var in continuous_vars:
            self.plotter.save_plot_to_ppt(
                lambda: sns.kdeplot(data=self.analyzer.df, x=var, hue=self.analyzer.target_variable, fill=True, common_norm=False),
                f'KDE Plot of {var} by {self.analyzer.target_variable}'
            )


def main():
    url = 'https://raw.githubusercontent.com/DwaipayanDutta/SAS_Data_model-Sample/main/Data/Updated_Data.py'
    
    # Initialize classes
    analyzer = DataAnalyzer(url)
    
    target_variable = input("Enter the target variable name (e.g., LI_FLAG): ")
    analyzer.set_target_variable(target_variable)

    ppt = Presentation()
    plotter = Plotter(ppt)
    
    eda = EDA(analyzer, plotter)

    eda.perform_kde_analysis()          # Perform KDE analysis on continuous variables.
    eda.perform_outlier_analysis()       # Perform outlier analysis.
    
    # Correlation Analysis
    correlation_matrix = analyzer.df[analyzer.get_continuous_vars()].corr()
    
    plt.figure(figsize=(12, 8))
    sns.heatmap(correlation_matrix, annot=True, fmt=".2f", cmap='coolwarm', square=True, cbar_kws={"shrink": .8})
    plt.title('Correlation Heatmap')
    plt.tight_layout()
    plt.savefig('temp_correlation_heatmap.png')
    plt.close()

    # Add heatmap to PowerPoint
    slide = ppt.slides.add_slide(ppt.slide_layouts[5]) 
    title = slide.shapes.title
    title.text = 'Correlation Heatmap'
    
    plotter.set_title_format(title)  # Set title formatting
    slide.shapes.add_picture('temp_correlation_heatmap.png', left=Inches(1), top=Inches(1.5), width=Inches(8))

    # Categorical Variables Analysis
    categorical_vars = analyzer.get_categorical_vars()

    for var in categorical_vars:
        plotter.save_plot_to_ppt(lambda: sns.countplot(x=analyzer.df[var], palette='Set2'), f'Count of {var}')

    # Bivariate Analysis: Continuous vs. Categorical with Box Plots
    for var in analyzer.get_continuous_vars():
        plotter.save_plot_to_ppt(lambda: sns.boxplot(x=analyzer.target_variable, y=var, data=analyzer.df, palette='pastel'), f'Box Plot of {var} by {analyzer.target_variable}')

    # Categorical vs. Categorical with Stacked Bar Plots
    for var in categorical_vars:
        if var != analyzer.target_variable:
            crosstab = pd.crosstab(analyzer.df[var], analyzer.df[analyzer.target_variable])
            crosstab.plot(kind='bar', stacked=True, figsize=(10, 5), color=sns.color_palette("Paired"))
            plt.title(f'Stacked Bar Plot of {var} by {analyzer.target_variable}')
            plt.tight_layout()
            plt.savefig('temp_stacked_bar.png')
            plt.close()

            # Add stacked bar plot to PowerPoint
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Title Only layout
            title = slide.shapes.title
            title.text = f'Stacked Bar Plot of {var} by {analyzer.target_variable}'
            plotter.set_title_format(title)  # Set title formatting
            
            slide.shapes.add_picture('temp_stacked_bar.png', left=Inches(1), top=Inches(1.5), width=Inches(8))

    # Save the PowerPoint presentation and print the path
    ppt_file_path = 'Data_Analysis_Presentation.pptx'
    ppt.save(ppt_file_path)
    
    print(f"Presentation saved at: {ppt_file_path}")


if __name__ == "__main__":
    main()